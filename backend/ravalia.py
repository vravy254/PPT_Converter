import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re

def change_text_properties(text_frame, font_name, font_size, font_color):
    for paragraph in text_frame.paragraphs:
        for r in paragraph.runs:
            if r.font.name != 'Arial': 
                r.font.name = font_name
            r.font.size = Pt(font_size)
            r.font.color.rgb = RGBColor(*font_color)

def html_email_check(text_frame, special_font):
    email_char = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
    http_char = r'https?://[^\s]+'
    
    for paragraph in text_frame.paragraphs:
        for r in paragraph.runs:
            if re.search(email_char, r.text) or re.search(http_char, r.text):
                r.font.name = special_font

def no_italics(text_frame):
    for paragraph in text_frame.paragraphs:
        for r in paragraph.runs:
            r.font.italic = False

def edit_presentation(ppt_path, output_path, title_font_name, title_font_size, title_font_color, text_font_name, text_font_size, text_font_color, text_special_font):
    file_path = Presentation(ppt_path)

    for slide in file_path.slides:
        if slide.shapes.title:
            html_email_check(slide.shapes.title.text_frame, text_special_font)
            change_text_properties(slide.shapes.title.text_frame, title_font_name, title_font_size, title_font_color)
            no_italics(slide.shapes.title.text_frame)

        for shape in slide.shapes:
            if not shape.has_text_frame or shape == slide.shapes.title:
                continue
            html_email_check(shape.text_frame, text_special_font)
            change_text_properties(shape.text_frame, text_font_name, text_font_size, text_font_color)
            no_italics(shape.text_frame)

    file_path.save(output_path)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python ravalia.py <input_ppt_path> <output_ppt_path>")
        sys.exit(1)

    ppt_path = sys.argv[1]
    output_path = sys.argv[2]

    # Hardcoded parameters for simplicity, can be extended to take inputs
    title_font_name = "72 Brand"
    title_font_size = 36
    title_font_color = (0, 0, 255)
    text_font_name = "72 Brand"
    text_font_size = 24
    text_font_color = (255, 0, 0)
    text_special_font = "Arial"

    edit_presentation(ppt_path, output_path, title_font_name, title_font_size, title_font_color, text_font_name, text_font_size, text_font_color, text_special_font)
    
    print(f"Presentation changed and saved as {output_path}")
